import csv
import json
import re
from pathlib import Path
from typing import Any, Iterable, List

from app.core.logger import node_log
from app.import_process.agent.state import ImportGraphState
from app.utils.path_util import PROJECT_ROOT
from app.utils.task_utils import add_done_task, add_running_task


SUPPORTED_CONVERT_EXTENSIONS = {
    ".txt",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
    ".html",
    ".htm",
    ".json",
}


def _read_text(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise UnicodeDecodeError("utf-8", data, 0, 1, "无法识别文本编码")


def _cell_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value).replace("\r", " ").replace("\n", " ").replace("|", "\\|").strip()


def _markdown_table(rows: Iterable[Iterable[Any]]) -> str:
    normalized = [[_cell_text(value) for value in row] for row in rows]
    normalized = [row for row in normalized if any(row)]
    if not normalized:
        return ""

    column_count = max(len(row) for row in normalized)
    normalized = [row + [""] * (column_count - len(row)) for row in normalized]
    header = normalized[0]
    if not any(header):
        header = [f"列{index}" for index in range(1, column_count + 1)]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * column_count) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in normalized[1:])
    return "\n".join(lines)


def _convert_docx(path: Path) -> str:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(path)
    output: List[str] = [f"# {path.stem}"]
    table_index = 0

    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style else ""
            heading = re.match(r"Heading\s*(\d+)", style_name, re.IGNORECASE)
            if heading:
                level = min(int(heading.group(1)) + 1, 6)
                output.append(f"{'#' * level} {text}")
            elif "list" in style_name.lower():
                output.append(f"- {text}")
            else:
                output.append(text)
        elif isinstance(child, CT_Tbl):
            table_index += 1
            table = Table(child, document)
            markdown = _markdown_table([[cell.text for cell in row.cells] for row in table.rows])
            if markdown:
                output.extend([f"## 表格 {table_index}", markdown])

    return "\n\n".join(output)


def _convert_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    output: List[str] = [f"# {path.stem}"]

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape and title_shape.has_text_frame else ""
        output.append(f"## 第 {slide_index} 页{f'：{title}' if title else ''}")

        for shape in sorted(slide.shapes, key=lambda item: (item.top, item.left)):
            if title_shape and shape.shape_id == title_shape.shape_id:
                continue
            if getattr(shape, "has_table", False):
                markdown = _markdown_table(
                    [[cell.text for cell in row.cells] for row in shape.table.rows]
                )
                if markdown:
                    output.append(markdown)
            elif getattr(shape, "has_text_frame", False):
                paragraphs = [paragraph.text.strip() for paragraph in shape.text_frame.paragraphs]
                output.extend(text for text in paragraphs if text)

    return "\n\n".join(output)


def _convert_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    output: List[str] = [f"# {path.stem}"]
    try:
        for worksheet in workbook.worksheets:
            output.append(f"## 工作表：{worksheet.title}")
            markdown = _markdown_table(worksheet.iter_rows(values_only=True))
            output.append(markdown or "（空工作表）")
    finally:
        workbook.close()
    return "\n\n".join(output)


def _convert_csv(path: Path) -> str:
    text = _read_text(path)
    rows = csv.reader(text.splitlines())
    return f"# {path.stem}\n\n{_markdown_table(rows)}"


def _convert_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_read_text(path), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    output: List[str] = []
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]):
        if tag.name != "table" and tag.find_parent("table"):
            continue
        if tag.name == "table":
            rows = [
                [cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])]
                for row in tag.find_all("tr")
            ]
            markdown = _markdown_table(rows)
            if markdown:
                output.append(markdown)
            continue

        text = tag.get_text(" ", strip=True)
        if not text:
            continue
        if tag.name.startswith("h"):
            output.append(f"{'#' * int(tag.name[1])} {text}")
        elif tag.name == "li":
            output.append(f"- {text}")
        else:
            output.append(text)

    return "\n\n".join(output) or soup.get_text("\n", strip=True)


def convert_file_to_markdown(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".txt":
        content = _read_text(path)
        return f"# {path.stem}\n\n{content}"
    if suffix == ".docx":
        return _convert_docx(path)
    if suffix == ".pptx":
        return _convert_pptx(path)
    if suffix == ".xlsx":
        return _convert_xlsx(path)
    if suffix == ".csv":
        return _convert_csv(path)
    if suffix in {".html", ".htm"}:
        return _convert_html(path)
    if suffix == ".json":
        data = json.loads(_read_text(path))
        return f"# {path.stem}\n\n```json\n{json.dumps(data, ensure_ascii=False, indent=2)}\n```"
    raise ValueError(f"暂不支持转换该文件类型：{suffix or '无扩展名'}")


@node_log("node_file_to_md")
def node_file_to_md(state: ImportGraphState) -> ImportGraphState:
    task_id = state["task_id"]
    add_running_task(task_id, "node_file_to_md")

    source_path = Path(state["local_file_path"])
    if not source_path.is_file():
        raise FileNotFoundError(f"待转换文件不存在：{source_path}")

    local_dir = Path(state.get("local_dir") or PROJECT_ROOT / "output")
    output_dir = local_dir / source_path.stem
    output_dir.mkdir(parents=True, exist_ok=True)
    markdown_path = output_dir / f"{source_path.stem}.md"
    markdown_content = convert_file_to_markdown(source_path).strip()
    if not markdown_content:
        raise ValueError(f"文件未解析出有效文本：{source_path.name}")

    markdown_path.write_text(markdown_content, encoding="utf-8")
    state["md_path"] = str(markdown_path)
    state["md_content"] = markdown_content
    add_done_task(task_id, "node_file_to_md")
    return state
